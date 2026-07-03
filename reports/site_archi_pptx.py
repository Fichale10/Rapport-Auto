"""
Génération d'un diagramme d'architecture de site (« qui porte qui »)
en PowerPoint **modifiable** : chaque nœud est une forme native (ovale),
chaque lien un connecteur, chaque libellé une zone de texte — tout reste
éditable dans PowerPoint.

Point d'entrée : ``build_architecture_pptx(data, site_name)`` où ``data`` est
la structure renvoyée par ``reports.views._build_site_architecture``.
"""

from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ── Palette (identique au rendu web) ──────────────────────────────────────────
C_ROOT   = RGBColor(0xFF, 0xC7, 0x2C)   # racine / IGW (jaune)
C_NODE   = RGBColor(0x3F, 0x7B, 0xD6)   # nœud porteur (bleu)
C_LEAF   = RGBColor(0x26, 0xA3, 0x5A)   # site terminal (vert)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_NAVY   = RGBColor(0x10, 0x2A, 0x6E)
C_DTEXT  = RGBColor(0x1C, 0x23, 0x40)
C_BLUE   = RGBColor(0x00, 0x30, 0x87)
C_RED    = RGBColor(0xFF, 0x5A, 0x5A)
C_EDGE   = RGBColor(0x9A, 0xA6, 0xC0)
C_SEC    = RGBColor(0xE0, 0xA8, 0x10)

# ── Dimensions de la diapo (16:9) ─────────────────────────────────────────────
SW = Inches(13.333)
SH = Inches(7.5)


def _compute_layout(data):
    """Replique l'algorithme de placement du rendu SVG (niveau + barycentre).

    Retourne (positions, level, children, ids) où ``positions[id] = (col, row)``
    en unités de grille (col = niveau, row = position verticale flottante)."""
    nodes = data.get('nodes', [])
    edges = data.get('edges', [])
    by_id = {n['id']: n for n in nodes}

    children, parents = {}, {}
    for e in edges:
        s, t = e.get('source'), e.get('target')
        if s in by_id and t in by_id:
            children.setdefault(s, []).append(t)
            parents.setdefault(t, []).append(s)

    # ── Niveau = plus long chemin depuis une racine ──
    level = {}

    def set_level(nid, lv, stack):
        if nid not in level or lv > level[nid]:
            level[nid] = lv
        else:
            return
        if stack.get(nid):
            return
        stack[nid] = True
        for c in children.get(nid, []):
            set_level(c, lv + 1, stack)
        stack[nid] = False

    for n in nodes:
        if not parents.get(n['id']):
            set_level(n['id'], 0, {})
    for n in nodes:
        level.setdefault(n['id'], 0)

    # ── Regroupement par niveau + ordonnancement barycentre ──
    by_level = {}
    for n in nodes:
        by_level.setdefault(level[n['id']], []).append(n['id'])
    levels = sorted(by_level.keys())
    row = {}
    for lv in levels:
        for i, nid in enumerate(by_level[lv]):
            row[nid] = i

    def bary(nid):
        ns = parents.get(nid, []) + children.get(nid, [])
        if not ns:
            return row[nid]
        return sum(row[x] for x in ns) / len(ns)

    for _ in range(5):
        for lv in levels:
            by_level[lv].sort(key=lambda a: (bary(a), a))
            for i, nid in enumerate(by_level[lv]):
                row[nid] = i

    max_rows = max((len(by_level[lv]) for lv in levels), default=1)
    max_level = levels[-1] if levels else 0

    positions = {}
    for nid in by_id:
        lv = level[nid]
        off = (max_rows - len(by_level[lv])) / 2.0   # centrage vertical par colonne
        positions[nid] = (lv, row[nid] + off)

    return positions, level, children, by_id, max_level, max_rows


def _set_line_dash(line, val='dash'):
    ln = line._get_or_add_ln()
    for pd in ln.findall(qn('a:prstDash')):
        ln.remove(pd)
    pd = ln.makeelement(qn('a:prstDash'), {'val': val})
    ln.append(pd)


def _set_line_arrow(line):
    ln = line._get_or_add_ln()
    for el in ln.findall(qn('a:tailEnd')):
        ln.remove(el)
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)


def _line(slide, x1, y1, x2, y2, color, w):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.shadow.inherit = False
    c.line.color.rgb = color
    c.line.width = Pt(w)
    return c


def _draw_tower(slide, cx, cy, d, color):
    """Dessine un pylône treillis (comme le rendu web) centré sur (cx, cy),
    à base de connecteurs natifs → reste modifiable dans PowerPoint."""
    top = cy - d * 0.24
    bot = cy + d * 0.20
    th = d * 0.09          # demi-largeur en haut
    bh = d * 0.19          # demi-largeur en bas
    lw = max(0.75, d * 2.0)

    def half(t):
        hw = bh + (th - bh) * t
        y = bot + (top - bot) * t
        return cx - hw, cx + hw, y

    # Montants
    _line(slide, cx - bh, bot, cx - th, top, color, lw)
    _line(slide, cx + bh, bot, cx + th, top, color, lw)

    # Traverses + croisillons
    steps = 3
    prev = None
    for i in range(steps + 1):
        l, r, y = half(i / steps)
        _line(slide, l, y, r, y, color, lw * 0.7)
        if prev:
            pl, pr, py = prev
            _line(slide, pl, py, r, y, color, lw * 0.5)
            _line(slide, pr, py, l, y, color, lw * 0.5)
        prev = (l, r, y)

    # Antenne
    _line(slide, cx, top, cx, top - d * 0.13, color, lw * 0.7)


def _add_architecture_slide(prs, data, site_name):
    """Ajoute à ``prs`` une diapo « architecture du site » (formes éditables)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # diapo vierge

    # ── Fond bleu foncé (comme le canvas web) ──
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    try:
        bg.fill.gradient()
        stops = bg.fill.gradient_stops
        stops[0].color.rgb = RGBColor(0x1C, 0x3F, 0x96)
        stops[1].color.rgb = RGBColor(0x0B, 0x1F, 0x55)
        try:
            bg.fill.gradient_angle = 55.0
        except Exception:
            pass
    except Exception:
        bg.fill.solid(); bg.fill.fore_color.rgb = C_NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False

    # ── En-tête ──
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.95))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = C_BLUE
    hdr.line.fill.background()
    hdr.shadow.inherit = False
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), Inches(0.95))
    band.fill.solid(); band.fill.fore_color.rgb = C_ROOT
    band.line.fill.background(); band.shadow.inherit = False

    tb = slide.shapes.add_textbox(Inches(0.35), Inches(0.14), Inches(12.5), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = '🏗️ Architecture du site'
    r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = C_WHITE
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = f'{site_name} — qui porte qui'
    r2.font.size = Pt(12); r2.font.color.rgb = C_ROOT

    positions, level, children, by_id, max_level, max_rows = _compute_layout(data)
    current = (data.get('current') or '').strip().upper()

    # ── Seuil SPOF (centralité) ──
    max_load = max((n.get('load', 0) for n in by_id.values()), default=0)
    spof_threshold = max(5, -(-max_load // 2))   # ceil(max_load / 2)

    # ── Zone de dessin ──
    draw_l, draw_t = 0.55, 1.25
    draw_w, draw_h = 12.25, 5.35

    # Diamètre des nœuds adapté au nombre de lignes
    node_d = 0.6
    if max_rows > 0:
        node_d = min(0.6, max(0.34, (draw_h / max_rows) * 0.42))

    # Étendue verticale réelle
    ys = [py for (_c, py) in positions.values()] or [0]
    min_y, max_y = min(ys), max(ys)
    span_y = (max_y - min_y) or 1

    def cx(col):
        if max_level <= 0:
            return draw_l + draw_w / 2.0
        return draw_l + node_d / 2.0 + col * (draw_w - node_d) / max_level

    def cy(py):
        if span_y <= 0:
            return draw_t + draw_h / 2.0
        return draw_t + node_d / 2.0 + (py - min_y) * (draw_h - node_d) / span_y

    # ── Connecteurs (dessinés avant les nœuds) ──
    C_EDGE_LIGHT = RGBColor(0xD5, 0xDD, 0xF0)
    for e in data.get('edges', []):
        s, t = e.get('source'), e.get('target')
        if s not in by_id or t not in by_id:
            continue
        x1, y1 = cx(positions[s][0]), cy(positions[s][1])
        x2, y2 = cx(positions[t][0]), cy(positions[t][1])
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.shadow.inherit = False
        line = conn.line
        sec = e.get('type') == 'secondary'
        trans = (e.get('trans') or '').upper()
        if sec:
            line.color.rgb = C_ROOT; line.width = Pt(2.0); _set_line_dash(line, 'dash')
        elif trans == 'FO':
            line.color.rgb = C_EDGE_LIGHT; line.width = Pt(3.2)
        elif trans == 'FH':
            line.color.rgb = C_EDGE_LIGHT; line.width = Pt(1.4); _set_line_dash(line, 'dash')
        elif trans == 'FTTM':
            line.color.rgb = C_EDGE_LIGHT; line.width = Pt(1.4)
        else:
            line.color.rgb = C_EDGE_LIGHT; line.width = Pt(2.2)
        # Orientation : flèche uniquement vers les sites terminaux (feuilles)
        if not children.get(t):
            _set_line_arrow(line)

    # ── Nœuds ──
    d = Inches(node_d)
    for nid, n in by_id.items():
        col, py = positions[nid]
        ccx, ccy = cx(col), cy(py)
        left = Inches(ccx - node_d / 2.0)
        top = Inches(ccy - node_d / 2.0)

        has_kids = bool(children.get(nid))
        is_root = level[nid] == 0
        is_current = nid == current
        is_spof = (n.get('load', 0) >= spof_threshold) and has_kids

        fill = C_ROOT if is_root else (C_NODE if has_kids else C_LEAF)

        # Anneau SPOF (rouge pointillé) sous le nœud
        if is_spof:
            ring = node_d * 0.32
            rg = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(ccx - node_d / 2.0 - ring / 2.0),
                Inches(ccy - node_d / 2.0 - ring / 2.0),
                Inches(node_d + ring), Inches(node_d + ring))
            rg.fill.background()
            rg.line.color.rgb = C_RED; rg.line.width = Pt(1.75)
            _set_line_dash(rg.line, 'sysDash'); rg.shadow.inherit = False

        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, d, d)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.shadow.inherit = False
        if is_current:
            shp.line.color.rgb = C_ROOT; shp.line.width = Pt(3.0)
        else:
            shp.line.color.rgb = C_WHITE; shp.line.width = Pt(1.75)

        # Icône pylône (treillis) dans le nœud
        _draw_tower(slide, ccx, ccy, node_d, C_NAVY if is_root else C_WHITE)

        # Libellé sous le nœud (texte blanc sur fond foncé)
        lbl = slide.shapes.add_textbox(
            Inches(ccx - 1.0), Inches(ccy + node_d / 2.0 + 0.03),
            Inches(2.0), Inches(0.3))
        ltf = lbl.text_frame; ltf.word_wrap = True
        lp = ltf.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run(); lr.text = n.get('name', '')
        lr.font.size = Pt(8.5)
        lr.font.bold = True
        lr.font.color.rgb = C_ROOT if is_current else C_WHITE

    # ── Légende (bas de diapo, texte blanc) ──
    legend = [
        (C_ROOT, 'Racine / IGW'),
        (C_NODE, 'Nœud porteur'),
        (C_LEAF, 'Site terminal'),
        (C_RED,  'Porteur critique (SPOF)'),
    ]
    lx = 0.55
    ly = 7.08
    for color, text in legend:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(lx), Inches(ly), Inches(0.16), Inches(0.16))
        dot.fill.solid(); dot.fill.fore_color.rgb = color
        dot.line.fill.background(); dot.shadow.inherit = False
        tbl = slide.shapes.add_textbox(Inches(lx + 0.20), Inches(ly - 0.05),
                                       Inches(2.6), Inches(0.28))
        tp = tbl.text_frame.paragraphs[0]
        tr = tp.add_run(); tr.text = text
        tr.font.size = Pt(9); tr.font.color.rgb = C_WHITE
        lx += max(1.75, 0.42 + len(text) * 0.082)


def build_architecture_pptx(data, site_name):
    """PPTX d'un seul site — retourne un ``BytesIO``."""
    return build_architectures_pptx([(data, site_name)])


def build_architectures_pptx(items):
    """PPTX multi-sites : ``items`` = liste de tuples ``(data, site_name)``.
    Une diapo par site. Retourne un ``BytesIO`` prêt au téléchargement."""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    for data, site_name in items:
        _add_architecture_slide(prs, data, site_name)
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

