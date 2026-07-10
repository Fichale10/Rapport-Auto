"""
Génération automatique du rapport PowerPoint — Comité Gestion des Incidents
"""

from io import BytesIO
from datetime import date

from pptx import Presentation

_MOIS_FR = [
    '', 'JANVIER', 'FEVRIER', 'MARS', 'AVRIL', 'MAI', 'JUIN',
    'JUILLET', 'AOUT', 'SEPTEMBRE', 'OCTOBRE', 'NOVEMBRE', 'DECEMBRE',
]

def _mois_label_fr(d):
    if not d:
        return ''
    return f'{_MOIS_FR[d.month]} {d.year}'
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette de couleurs ───────────────────────────────────────────────────────
C_BLUE   = RGBColor(0x00, 0x30, 0x87)
C_BLUE2  = RGBColor(0x00, 0x47, 0xCC)
C_BLUE3  = RGBColor(0x1E, 0x3A, 0x6E)
C_YELL   = RGBColor(0xFF, 0xC7, 0x2C)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY  = RGBColor(0xF8, 0xFA, 0xFF)
C_MGRAY  = RGBColor(0xE8, 0xED, 0xF5)
C_DTEXT  = RGBColor(0x33, 0x33, 0x44)
C_BLUE_T = RGBColor(0x2A, 0x4A, 0x80)

C_GREEN_BG  = RGBColor(0xC6, 0xEF, 0xCE)
C_GREEN_FG  = RGBColor(0x27, 0x62, 0x21)
C_YELL_BG   = RGBColor(0xFF, 0xEB, 0x9C)
C_YELL_FG   = RGBColor(0x9C, 0x65, 0x00)
C_RED_BG    = RGBColor(0xFF, 0xC7, 0xCE)
C_RED_FG    = RGBColor(0x9C, 0x00, 0x06)

# ── Dimensions ────────────────────────────────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)
HDR_H = Inches(1.08)
MARGIN = Inches(0.35)
CONTENT_TOP = Inches(1.18)
CONTENT_H   = SH - CONTENT_TOP - Inches(0.25)


# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS BAS NIVEAU
# ═══════════════════════════════════════════════════════════════════════════════

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, rgb=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if rgb:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb
    else:
        shp.fill.background()
    shp.line.fill.background()
    return shp


def _txt(slide, text, l, t, w, h, size=11, bold=False, color=C_DTEXT,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def _fmt(seconds):
    if not seconds:
        return '0:00:00'
    s = int(seconds)
    return f'{s//3600}:{(s%3600)//60:02d}:{s%60:02d}'


# ═══════════════════════════════════════════════════════════════════════════════
# COMPOSANTS SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def _header(slide, title, subtitle='', mois_label=''):
    _rect(slide, 0, 0, SW, HDR_H, C_BLUE)
    _rect(slide, 0, 0, Inches(0.07), HDR_H, C_YELL)    # bande jaune gauche
    tag = 'COMITE GESTION DES INCIDENTS'
    _txt(slide, tag, MARGIN, Inches(0.06), Inches(10), Inches(0.28),
         size=8, bold=True, color=C_YELL)
    _txt(slide, title, MARGIN, Inches(0.30), Inches(10), Inches(0.50),
         size=17, bold=True, color=C_WHITE)
    if subtitle:
        _txt(slide, subtitle, MARGIN, Inches(0.78), Inches(9), Inches(0.25),
             size=9, color=RGBColor(0xB0, 0xC4, 0xE8))
    if mois_label:
        _txt(slide, mois_label, Inches(10.5), Inches(0.32), Inches(2.5), Inches(0.4),
             size=13, bold=True, color=C_YELL, align=PP_ALIGN.RIGHT)


def _table(slide, headers, rows,
           left=MARGIN, top=CONTENT_TOP, width=None, height=None,
           col_widths=None, hdr_bg=C_BLUE3, alt=True, cell_fmts=None,
           font_size=9, hdr_size=9):
    if width  is None: width  = SW - 2 * MARGIN
    if height is None: height = CONTENT_H - Inches(0.1)

    nc = len(headers)
    nr = len(rows) + 1
    tbl = slide.shapes.add_table(nr, nc, left, top, width, height).table

    # Largeurs colonnes
    if col_widths:
        tot = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(width * cw / tot)
    else:
        cw = width // nc
        for i in range(nc):
            tbl.columns[i].width = cw

    # Hauteurs lignes
    rh = height // nr
    for i in range(nr):
        tbl.rows[i].height = rh

    def _cell(cell, val, bg=None, fg=C_DTEXT, bold=False,
              align=PP_ALIGN.CENTER, fs=9):
        # Clear existing text
        cell.text = ''
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        else:
            try:
                cell.fill.background()
            except Exception:
                pass
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(val) if val is not None else ''
        r.font.size = Pt(fs)
        r.font.bold = bold
        r.font.color.rgb = fg

    # Entête
    for j, h in enumerate(headers):
        _cell(tbl.cell(0, j), h, bg=hdr_bg, fg=C_WHITE, bold=True, fs=hdr_size)

    # Données
    for i, row in enumerate(rows):
        bg_row = C_LGRAY if (alt and i % 2 == 1) else None
        for j, val in enumerate(row):
            fmt = cell_fmts.get((i, j)) if cell_fmts else None
            if fmt:
                cb, cf = fmt
                _cell(tbl.cell(i+1, j), val, bg=cb, fg=cf, bold=True, fs=font_size)
            else:
                is_lbl = (j == 0)
                _cell(tbl.cell(i+1, j), val, bg=bg_row,
                      fg=C_BLUE3 if is_lbl else C_DTEXT,
                      bold=is_lbl,
                      align=PP_ALIGN.LEFT if is_lbl else PP_ALIGN.CENTER,
                      fs=font_size)

    return tbl


def _kpi_bar(slide, kpis):
    """kpis = list of (label, value, unit, color)"""
    n = len(kpis)
    w = (SW - 2 * MARGIN) // n
    for i, (lbl, val, unit, color) in enumerate(kpis):
        x = MARGIN + i * w
        y = CONTENT_TOP
        _rect(slide, x + Inches(0.06), y, w - Inches(0.12), Inches(1.6), color)
        _txt(slide, str(val), x + Inches(0.06), y + Inches(0.25), w - Inches(0.12), Inches(0.8),
             size=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, lbl, x + Inches(0.06), y + Inches(1.1), w - Inches(0.12), Inches(0.35),
             size=10, bold=True, color=RGBColor(0xE8, 0xF0, 0xFF), align=PP_ALIGN.CENTER)
        if unit:
            _txt(slide, unit, x + Inches(0.06), y + Inches(0.08), w - Inches(0.12), Inches(0.25),
                 size=8, color=RGBColor(0xC0, 0xD4, 0xF0), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES FIXES
# ═══════════════════════════════════════════════════════════════════════════════

def _cover(prs, mois_label, generated_on):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, C_BLUE)
    _rect(sl, 0, 0, SW, Inches(0.12), C_YELL)
    _rect(sl, 0, SH - Inches(0.12), SW, Inches(0.12), C_YELL)
    # Bandes décoratives
    _rect(sl, SW - Inches(0.12), 0, Inches(0.12), SH, C_YELL)
    _rect(sl, 0, 0, Inches(0.12), SH, C_YELL)

    _txt(sl, 'Yas Togo / DT / DOC / iSOC', MARGIN, Inches(0.5), SW - 2*MARGIN, Inches(0.4),
         size=11, color=RGBColor(0x80, 0xA0, 0xD0), align=PP_ALIGN.CENTER)
    _txt(sl, 'COMITÉ', MARGIN, Inches(1.5), SW - 2*MARGIN, Inches(1.1),
         size=60, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, 'GESTION DES INCIDENTS', MARGIN, Inches(2.5), SW - 2*MARGIN, Inches(0.9),
         size=40, bold=True, color=C_YELL, align=PP_ALIGN.CENTER)
    _txt(sl, mois_label.upper(), MARGIN, Inches(3.7), SW - 2*MARGIN, Inches(0.8),
         size=32, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, f'Généré le {generated_on}', MARGIN, SH - Inches(0.8),
         SW - 2*MARGIN, Inches(0.35),
         size=10, color=RGBColor(0x80, 0xA0, 0xD0), align=PP_ALIGN.CENTER)


def _section(prs, num, title, icon=''):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, C_BLUE)
    _rect(sl, 0, 0, Inches(0.12), SH, C_YELL)
    _rect(sl, SW - Inches(0.12), 0, Inches(0.12), SH, C_YELL)
    _txt(sl, str(num).zfill(2), MARGIN + Inches(0.5), Inches(1.5), Inches(3), Inches(3.5),
         size=140, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT)
    _txt(sl, icon, Inches(5), Inches(2.0), Inches(2), Inches(2),
         size=72, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, title, MARGIN, Inches(5.5), SW - 2*MARGIN, Inches(1.2),
         size=36, bold=True, color=C_YELL, align=PP_ALIGN.CENTER)


def _closing(prs):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, C_BLUE)
    _rect(sl, 0, 0, SW, Inches(0.12), C_YELL)
    _rect(sl, 0, SH - Inches(0.12), SW, Inches(0.12), C_YELL)
    _txt(sl, 'MERCI', MARGIN, Inches(2.5), SW - 2*MARGIN, Inches(2),
         size=72, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, 'Yas Togo / DT / DOC / iSOC / GDI',
         MARGIN, SH - Inches(0.8), SW - 2*MARGIN, Inches(0.35),
         size=10, color=RGBColor(0x80, 0xA0, 0xD0), align=PP_ALIGN.CENTER)


def _def_slide(prs, mois_label):
    sl = _blank(prs)
    _header(sl, 'Définition DR1 / DR2', mois_label=mois_label)
    headers = ['Code', 'Indicateur', 'Définition', 'Seuil (2G/3G/4G)']
    rows = [
        ['DR1', "Nombre d'indisponibilités d'une station",
         "Nombre de fois qu'une même station est indisponible sur le mois",
         '≤ 2'],
        ['DR2', "Délai d'indisponibilité d'une station",
         "Délai d'indisponibilité par jour d'une station > 3h constitue un DR2",
         '≤ 3H'],
    ]
    _table(sl, headers, rows,
           col_widths=[1, 3, 6, 2],
           top=CONTENT_TOP + Inches(0.8),
           height=Inches(2.2),
           font_size=11, hdr_size=11)
    _txt(sl, '• Un DR1 = site indisponible plus de 2 fois dans le mois',
         MARGIN, CONTENT_TOP + Inches(3.5), SW - 2*MARGIN, Inches(0.4),
         size=12, color=C_BLUE3)
    _txt(sl, '• Un DR2 = une indisponibilité dont la durée dépasse 3 heures',
         MARGIN, CONTENT_TOP + Inches(4.0), SW - 2*MARGIN, Inches(0.4),
         size=12, color=C_BLUE3)


# ═══════════════════════════════════════════════════════════════════════════════
# FETCH DATA
# ═══════════════════════════════════════════════════════════════════════════════

def _fetch_mobile(mois):
    from .models import Incident
    from django.db.models import Count, Sum

    qs = Incident.objects.filter(domain='mobile')
    if mois:
        qs = qs.filter(mois_rapport=mois)

    total = qs.count()

    # DR1 violations: sites avec > 2 incidents
    dr1_rows = list(
        qs.exclude(site_name='')
        .values('site_name', 'region', 'cause')
        .annotate(cnt=Count('id'))
        .filter(cnt__gt=2)
        .order_by('-cnt')[:15]
    )

    # DR2 data
    dr2_qs = Incident.objects.filter(domain='dr2')
    dr2_mois = dr2_qs.order_by('-mois_rapport').values_list('mois_rapport', flat=True).first()
    if dr2_mois:
        dr2_qs = dr2_qs.filter(mois_rapport=dr2_mois)
    total_dr2 = dr2_qs.count()

    dr2_by_region = {r['region'].upper(): r['cnt'] for r in
                     dr2_qs.exclude(region='').values('region').annotate(cnt=Count('id'))}
    dr2_by_escalade = {e['escalade']: e['cnt'] for e in
                       dr2_qs.exclude(escalade='').values('escalade').annotate(cnt=Count('id'))}

    # Stats par région
    region_stats = []
    for r in qs.exclude(region='').values('region').annotate(
        nb=Count('id'), dur=Sum('duration_sec')
    ).order_by('-nb'):
        reg = r['region'].upper()
        nb = r['nb']
        dur = r['dur'] or 0
        mttr_sec = dur / nb if nb else 0
        dr2_r = dr2_by_region.get(reg, 0)
        eff = round((nb - dr2_r) / nb * 100) if nb else 0
        region_stats.append({
            'region': r['region'], 'nb': nb,
            'mttr': _fmt(mttr_sec), 'dr2': dr2_r, 'eff': eff,
        })

    # Top sites (DR1)
    top_sites = list(
        qs.exclude(site_name='').values('site_name', 'region')
        .annotate(cnt=Count('id')).order_by('-cnt')[:15]
    )

    # Points bloquants
    pb_rows = list(
        qs.exclude(point_bloquant='').values('point_bloquant')
        .annotate(cnt=Count('id')).order_by('-cnt')[:10]
    )

    # Top causes
    cause_rows = list(
        qs.exclude(cause='').values('cause')
        .annotate(cnt=Count('id')).order_by('-cnt')[:10]
    )

    # DR2 par escalade
    esc_rows = list(
        dr2_qs.exclude(escalade='').values('escalade')
        .annotate(nb=Count('id'), dur=Sum('duration_sec')).order_by('-nb')
    )
    for e in esc_rows:
        nb = e['nb']
        dur = e['dur'] or 0
        e['mttr'] = _fmt(dur / nb if nb else 0)
        e['outage'] = _fmt(dur)
        e['pct'] = f"{round(nb / total_dr2 * 100)}%" if total_dr2 else '0%'

    return {
        'total': total, 'total_dr2': total_dr2,
        'tdr1': len(dr1_rows), 'dr1_rows': dr1_rows,
        'region_stats': region_stats, 'top_sites': top_sites,
        'pb_rows': pb_rows, 'cause_rows': cause_rows,
        'esc_rows': esc_rows,
    }


def _fetch_domain(domain, mois):
    from .models import Incident
    from django.db.models import Count, Sum

    qs = Incident.objects.filter(domain=domain)
    if mois:
        qs = qs.filter(mois_rapport=mois)
    if not qs.exists():
        qs = Incident.objects.filter(domain=domain)

    total = qs.count()
    total_dur = qs.aggregate(d=Sum('duration_sec'))['d'] or 0

    by_region = list(
        qs.exclude(region='').values('region')
        .annotate(nb=Count('id'), dur=Sum('duration_sec')).order_by('-nb')
    )
    for r in by_region:
        nb = r['nb']
        r['mttr'] = _fmt((r['dur'] or 0) / nb if nb else 0)
        r['outage'] = _fmt(r['dur'] or 0)

    by_escalade = list(
        qs.exclude(escalade='').values('escalade')
        .annotate(nb=Count('id'), dur=Sum('duration_sec')).order_by('-nb')
    )
    for e in by_escalade:
        nb = e['nb']
        e['mttr']   = _fmt((e['dur'] or 0) / nb if nb else 0)
        e['outage'] = _fmt(e['dur'] or 0)

    by_nature = list(
        qs.exclude(nature='').values('nature')
        .annotate(nb=Count('id')).order_by('-nb')[:10]
    )

    incidents = list(qs.exclude(nature='').values(
        'nature', 'region', 'escalade', 'duration_sec', 'site_name', 'cause', 'status'
    ).order_by('-duration_sec')[:20])
    for inc in incidents:
        inc['dur_fmt'] = _fmt(inc['duration_sec'])

    return {
        'total': total, 'total_dur': _fmt(total_dur),
        'by_region': by_region, 'by_escalade': by_escalade,
        'by_nature': by_nature, 'incidents': incidents,
    }


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES RÉSEAU MOBILE
# ═══════════════════════════════════════════════════════════════════════════════

def _slide_dr1(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, f'Cas de violation DR1   —   TDR1 = {d["tdr1"]} sites', mois_label=mois_label)
    rows = d['dr1_rows'][:15]
    tbl_rows = [[r['site_name'], r['cnt'], r['region'],
                 str(r.get('cause', '') or '')[:45]] for r in rows]
    fmts = {}
    for i, r in enumerate(rows):
        if r['cnt'] >= 10:
            fmts[(i, 1)] = (C_RED_BG, C_RED_FG)
        elif r['cnt'] >= 5:
            fmts[(i, 1)] = (C_YELL_BG, C_YELL_FG)
        else:
            fmts[(i, 1)] = (C_GREEN_BG, C_GREEN_FG)
    _table(sl, ['SITE', 'COUNT', 'RÉGION', 'CAUSE'],
           tbl_rows, col_widths=[3, 1, 2, 6],
           cell_fmts=fmts, font_size=9)


def _slide_dr2_overview(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, f'DR2 Trend   —   TDR2 = {d["total_dr2"]} DR2', mois_label=mois_label)
    kpis = [
        ('Total Incidents', d['total'], 'RÉSEAU MOBILE', C_BLUE),
        ('Total DR2', d['total_dr2'], 'VIOLATIONS > 3H', C_RED_FG),
        ('Violations DR1', d['tdr1'], 'SITES > 2 FOIS', RGBColor(0xD9, 0x77, 0x06)),
        ('Efficacité glob.', f"{round((d['total']-d['total_dr2'])/d['total']*100) if d['total'] else 0}%",
         'RÉSOLUTION < 3H', RGBColor(0x05, 0x96, 0x69)),
    ]
    _kpi_bar(sl, kpis)

    # Table par région
    _txt(sl, 'Efficacité DR2 par Région', MARGIN,
         CONTENT_TOP + Inches(1.8), SW - 2*MARGIN, Inches(0.35),
         size=12, bold=True, color=C_BLUE3)
    rows = [[r['region'], r['nb'], r['mttr'], r['dr2'],
             f"{r['eff']}%"] for r in d['region_stats']]
    fmts = {}
    for i, r in enumerate(d['region_stats']):
        eff = r['eff']
        if eff >= 90:
            fmts[(i, 4)] = (C_GREEN_BG, C_GREEN_FG)
        elif eff >= 80:
            fmts[(i, 4)] = (C_YELL_BG, C_YELL_FG)
        else:
            fmts[(i, 4)] = (C_RED_BG, C_RED_FG)
    _table(sl, ['RÉGION', 'Nb INCIDENTS', 'MTTR', 'Nb DR2', '% EFFICACITÉ'],
           rows, col_widths=[3, 2, 2, 2, 2],
           top=CONTENT_TOP + Inches(2.2),
           height=Inches(2.8),
           cell_fmts=fmts, font_size=10)


def _slide_dr2_metier(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, 'Violation DR2 / Métiers', mois_label=mois_label)
    rows = [[e['escalade'], e['nb'], e['mttr'], e['outage'], e['pct']]
            for e in d['esc_rows']]
    _table(sl, ['ESCALADE / MÉTIER', 'Nb DR2', 'MTTR', 'OUTAGE TOTAL', '% du Total'],
           rows, col_widths=[4, 2, 2, 2, 2], font_size=10)


def _slide_top_sites(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, 'Top Sites Récurrents', mois_label=mois_label)
    rows = [[r['site_name'], r['region'], r['cnt']] for r in d['top_sites'][:15]]
    fmts = {}
    for i, r in enumerate(d['top_sites'][:15]):
        if r['cnt'] >= 15:
            fmts[(i, 2)] = (C_RED_BG, C_RED_FG)
        elif r['cnt'] >= 8:
            fmts[(i, 2)] = (C_YELL_BG, C_YELL_FG)
    _table(sl, ['SITE', 'RÉGION', 'NB INCIDENTS'],
           rows, col_widths=[5, 3, 2], cell_fmts=fmts, font_size=10)


def _slide_causes(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, 'Top Causes des Incidents', mois_label=mois_label)
    total = d['total'] or 1
    rows = [[r['cause'][:60], r['cnt'],
             f"{round(r['cnt']/total*100)}%"] for r in d['cause_rows']]
    _table(sl, ['CAUSE', 'NB INCIDENTS', '% du Total'],
           rows, col_widths=[7, 2, 2], font_size=10)


def _slide_points_bloquants(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, 'Points Bloquants', mois_label=mois_label)
    total_pb = sum(r['cnt'] for r in d['pb_rows']) or 1
    rows = [[r['point_bloquant'][:55], r['cnt'],
             f"{round(r['cnt']/total_pb*100)}%"] for r in d['pb_rows']]
    # Ligne total
    rows.append(['TOTAL', sum(r['cnt'] for r in d['pb_rows']), '100%'])
    fmts = {(len(rows)-1, 0): (C_BLUE, C_WHITE),
            (len(rows)-1, 1): (C_BLUE, C_WHITE),
            (len(rows)-1, 2): (C_BLUE, C_WHITE)}
    _table(sl, ['CAUSE', 'COUNT', '% du Total'],
           rows, col_widths=[7, 2, 2], cell_fmts=fmts, font_size=10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES RÉSEAU FIXE
# ═══════════════════════════════════════════════════════════════════════════════

def _slide_fixe(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, f'Réseau Fixe — {d["total"]} incidents', mois_label=mois_label)
    kpis = [
        ('Total Incidents', d['total'], 'RÉSEAU FIXE', C_BLUE2),
        ('Outage Total', d['total_dur'], '', RGBColor(0xD9, 0x77, 0x06)),
    ]
    _kpi_bar(sl, kpis)

    # Types par nature (résumé)
    if d['by_nature']:
        _txt(sl, 'Types d\'incidents', MARGIN,
             CONTENT_TOP + Inches(1.8), SW / 2 - MARGIN - Inches(0.1), Inches(0.35),
             size=11, bold=True, color=C_BLUE3)
        nat_rows = [[r['nature'][:40], r['nb']] for r in d['by_nature'][:8]]
        _table(sl, ['NATURE', 'Nb'],
               nat_rows, col_widths=[7, 1],
               left=MARGIN, top=CONTENT_TOP + Inches(2.2),
               width=SW // 2 - MARGIN - Inches(0.1),
               height=Inches(3.0), font_size=9)

    # Par région
    if d['by_region']:
        _txt(sl, 'Par Région', SW // 2 + Inches(0.1),
             CONTENT_TOP + Inches(1.8), SW // 2 - MARGIN - Inches(0.1), Inches(0.35),
             size=11, bold=True, color=C_BLUE3)
        reg_rows = [[r['region'], r['nb'], r['mttr']] for r in d['by_region']]
        _table(sl, ['RÉGION', 'Nb', 'MTTR'],
               reg_rows, col_widths=[3, 1, 2],
               left=SW // 2 + Inches(0.1),
               top=CONTENT_TOP + Inches(2.2),
               width=SW // 2 - MARGIN - Inches(0.1),
               height=Inches(3.0), font_size=9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES TRANSPORT
# ═══════════════════════════════════════════════════════════════════════════════

def _slide_transport(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, f'Transport — {d["total"]} incidents', mois_label=mois_label)
    kpis = [
        ('Total Incidents', d['total'], 'TRANSPORT', C_BLUE2),
        ('Outage Total', d['total_dur'], '', RGBColor(0xD9, 0x77, 0x06)),
    ]
    _kpi_bar(sl, kpis)

    # Par escalade
    if d['by_escalade']:
        _txt(sl, 'Par Métier', MARGIN,
             CONTENT_TOP + Inches(1.8), SW // 2 - MARGIN - Inches(0.1), Inches(0.35),
             size=11, bold=True, color=C_BLUE3)
        esc_rows = [[e['escalade'][:30], e['nb'], e['mttr'], e['outage']]
                    for e in d['by_escalade']]
        _table(sl, ['MÉTIER', 'Nb', 'MTTR', 'OUTAGE'],
               esc_rows, col_widths=[4, 1, 2, 2],
               left=MARGIN, top=CONTENT_TOP + Inches(2.2),
               width=SW // 2 - MARGIN - Inches(0.1),
               height=Inches(3.2), font_size=9)

    # Par région
    if d['by_region']:
        _txt(sl, 'Par Région', SW // 2 + Inches(0.1),
             CONTENT_TOP + Inches(1.8), SW // 2 - MARGIN - Inches(0.1), Inches(0.35),
             size=11, bold=True, color=C_BLUE3)
        reg_rows = [[r['region'], r['nb'], r['mttr'], r['outage']]
                    for r in d['by_region']]
        _table(sl, ['RÉGION', 'Nb', 'MTTR', 'OUTAGE'],
               reg_rows, col_widths=[3, 1, 2, 2],
               left=SW // 2 + Inches(0.1),
               top=CONTENT_TOP + Inches(2.2),
               width=SW // 2 - MARGIN - Inches(0.1),
               height=Inches(3.2), font_size=9)


def _slide_transport_detail(prs, d, mois_label):
    """Slide incidents transport avec impact."""
    sl = _blank(prs)
    _header(sl, 'Détails Incidents Transport', mois_label=mois_label)
    incs = d['incidents'][:15]
    rows = [[i['site_name'][:25] or '—', i['region'],
             i['escalade'][:20], i['dur_fmt'],
             (i['cause'] or '')[:35]] for i in incs]
    _table(sl, ['SITE / LIEN', 'RÉGION', 'MÉTIER', 'DURÉE', 'CAUSE'],
           rows, col_widths=[3, 2, 2, 2, 4], font_size=9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES IGW
# ═══════════════════════════════════════════════════════════════════════════════

def _slide_igw(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, f'IGW — Disponibilité et Trafic   ({d["total"]} incidents)', mois_label=mois_label)
    kpis = [
        ('Total Incidents', d['total'], 'IGW', C_BLUE2),
        ('Outage Total', d['total_dur'], '', RGBColor(0xD9, 0x77, 0x06)),
    ]
    _kpi_bar(sl, kpis)

    if d['incidents']:
        _txt(sl, 'Incidents critiques IGW', MARGIN,
             CONTENT_TOP + Inches(1.8), SW - 2*MARGIN, Inches(0.35),
             size=11, bold=True, color=C_BLUE3)
        incs = d['incidents'][:12]
        rows = [[i['site_name'][:30] or i['nature'][:30],
                 i['escalade'][:20],
                 (i['cause'] or '')[:35],
                 i['dur_fmt']] for i in incs]
        fmts = {}
        for i, inc in enumerate(incs):
            if (inc['duration_sec'] or 0) > 72*3600:
                fmts[(i, 3)] = (C_RED_BG, C_RED_FG)
            elif (inc['duration_sec'] or 0) > 24*3600:
                fmts[(i, 3)] = (C_YELL_BG, C_YELL_FG)
        _table(sl, ['LIEN / SITE', 'ESCALADE', 'CAUSE', 'DURÉE'],
               rows, col_widths=[3, 2, 5, 2],
               top=CONTENT_TOP + Inches(2.2),
               height=Inches(3.0),
               cell_fmts=fmts, font_size=9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES CORE
# ═══════════════════════════════════════════════════════════════════════════════

def _slide_core(prs, d, mois_label):
    sl = _blank(prs)
    _header(sl, f'Core — {d["total"]} incidents', mois_label=mois_label)
    if d['incidents']:
        incs = d['incidents'][:10]
        rows = [[i['nature'][:45] or '—',
                 (i['cause'] or '')[:35],
                 i['escalade'][:20],
                 i['dur_fmt']] for i in incs]
        fmts = {}
        for i, inc in enumerate(incs):
            if (inc['duration_sec'] or 0) > 72*3600:
                fmts[(i, 3)] = (C_RED_BG, C_RED_FG)
        _table(sl, ['NATURE INCIDENT', 'CAUSE', 'ESCALADE', 'DURÉE'],
               rows, col_widths=[5, 4, 2, 2], cell_fmts=fmts, font_size=10)
    else:
        _txt(sl, 'Aucun incident enregistré pour cette période.',
             MARGIN, CONTENT_TOP + Inches(1), SW - 2*MARGIN, Inches(0.5),
             size=14, color=C_BLUE3)


# ═══════════════════════════════════════════════════════════════════════════════
# POINT D'ENTRÉE PRINCIPAL
# ═══════════════════════════════════════════════════════════════════════════════

def generate_report(mois_mobile=None, mois_dr2=None, mois_fixe=None,
                    mois_transport=None, mois_igw=None, mois_core=None,
                    generated_on=''):
    """
    Génère le rapport PPTX complet.
    Retourne un BytesIO prêt à être envoyé en réponse HTTP.
    """
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    # Label de mois pour l'affichage
    if mois_mobile:
        mois_label = _mois_label_fr(mois_mobile)
    else:
        from datetime import date as _d
        mois_label = _mois_label_fr(_d.today())

    if not generated_on:
        from datetime import date as _d
        generated_on = _d.today().strftime('%d/%m/%Y')

    # ── Couverture ──────────────────────────────────────────────────────────────
    _cover(prs, mois_label, generated_on)

    # ── Section 01 : Réseau Mobile ──────────────────────────────────────────────
    mob = _fetch_mobile(mois_mobile)

    _section(prs, 1, 'RÉSEAU MOBILE', '📡')
    _def_slide(prs, mois_label)
    _slide_dr1(prs, mob, mois_label)
    _slide_dr2_overview(prs, mob, mois_label)
    _slide_dr2_metier(prs, mob, mois_label)
    _slide_top_sites(prs, mob, mois_label)
    _slide_causes(prs, mob, mois_label)
    _slide_points_bloquants(prs, mob, mois_label)

    # ── Section 02 : Réseau Fixe ────────────────────────────────────────────────
    fixe = _fetch_domain('fixe', mois_fixe)
    _section(prs, 2, 'RÉSEAU FIXE', '☎️')
    _slide_fixe(prs, fixe, mois_label)

    # ── Section 03 : Transport ──────────────────────────────────────────────────
    transport = _fetch_domain('transport', mois_transport)
    _section(prs, 3, 'TRANSPORT', '🔗')
    _slide_transport(prs, transport, mois_label)
    if transport['incidents']:
        _slide_transport_detail(prs, transport, mois_label)

    # ── Section 04 : IGW ────────────────────────────────────────────────────────
    igw = _fetch_domain('igw', mois_igw)
    _section(prs, 4, 'IGW', '🔌')
    _slide_igw(prs, igw, mois_label)

    # ── Section 05 : Core ───────────────────────────────────────────────────────
    core = _fetch_domain('core', mois_core)
    _section(prs, 5, 'CORE', '🌐')
    _slide_core(prs, core, mois_label)

    # ── Fermeture ───────────────────────────────────────────────────────────────
    _closing(prs)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# RAPPORT GDI CORE — « Disponibilité et trafic IGW » (1 slide, depuis upload)
# ═══════════════════════════════════════════════════════════════════════════════

def _gdi_core_slide(prs, table_rows, period_label, page_label=''):
    """Construit une slide « Incidents core » avec les lignes fournies."""
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, C_WHITE)

    # Titres
    _txt(sl, 'COMITE GESTION DES INCIDENTS', MARGIN, Inches(0.18),
         Inches(9), Inches(0.45), size=22, bold=True, color=C_BLUE)
    _txt(sl, 'Disponibilité et trafic IGW', MARGIN, Inches(0.62),
         Inches(9), Inches(0.40), size=16, bold=True,
         color=RGBColor(0xE3, 0x00, 0x13))
    if period_label:
        _txt(sl, period_label, Inches(9.0), Inches(0.22), Inches(2.6), Inches(0.4),
             size=14, bold=True, color=C_BLUE, align=PP_ALIGN.RIGHT)

    # Logo Yas (haut-droite)
    try:
        from .gdi_core import yas_logo_bytes
        logo = yas_logo_bytes()
        if logo:
            from PIL import Image as _PILImage
            _img = _PILImage.open(BytesIO(logo))
            ratio = _img.height / _img.width
            lw = Inches(1.25)
            lh = Inches(1.25 * ratio)
            sl.shapes.add_picture(BytesIO(logo), SW - MARGIN - lw, Inches(0.18), lw, lh)
    except Exception:
        pass

    # Étiquette « Incidents core »
    _rect(sl, MARGIN, Inches(1.20), Inches(2.0), Inches(0.46), C_YELL)
    _txt(sl, 'Incidents core', MARGIN, Inches(1.24), Inches(2.0), Inches(0.40),
         size=14, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
    if page_label:
        _txt(sl, page_label, SW - MARGIN - Inches(2.4), Inches(1.24),
             Inches(2.4), Inches(0.40), size=12, bold=True, color=C_BLUE3,
             align=PP_ALIGN.RIGHT)

    headers = ["Nature de l'incident", 'Impact - Service', 'Cause', 'Escalade', 'Duration']

    if not table_rows:
        _txt(sl, 'Aucun incident dans le fichier importé.',
             MARGIN, Inches(2.2), SW - 2*MARGIN, Inches(0.5),
             size=14, color=C_BLUE3)
        return

    nr = len(table_rows) + 1
    top = Inches(1.95)
    height = SH - top - Inches(0.45)
    tbl = sl.shapes.add_table(
        nr, 5, MARGIN, top, SW - 2*MARGIN, height).table
    tbl.first_row = False
    tbl.horz_banding = False

    widths = [3.3, 2.5, 2.9, 2.2, 1.5]
    tw = SW - 2*MARGIN
    tot = sum(widths)
    for i, cw in enumerate(widths):
        tbl.columns[i].width = int(tw * cw / tot)

    # En-tête plus fine, lignes de données hautes (style Image 2)
    hdr_h = Inches(0.5)
    tbl.rows[0].height = hdr_h
    body_h = int((height - hdr_h) / len(table_rows))
    for i in range(1, nr):
        tbl.rows[i].height = body_h

    from pptx.enum.text import MSO_ANCHOR

    def _c(cell, val, bg, fg, bold, align, fs=11):
        cell.text = ''
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(val) if val is not None else ''
        r.font.size = Pt(fs)
        r.font.bold = bold
        r.font.color.rgb = fg

    # En-tête jaune, texte bleu (gras, aligné à gauche)
    for j, h in enumerate(headers):
        _c(tbl.cell(0, j), h, C_YELL, C_BLUE, True, PP_ALIGN.LEFT, fs=12)

    C_NAVY_CELL = RGBColor(0x0D, 0x24, 0x61)
    C_GRAY_CELL = RGBColor(0xE8, 0xEC, 0xF4)
    # Réduire la police quand beaucoup de lignes pour rester lisible
    fs = 11 if len(table_rows) <= 9 else 9
    for i, row in enumerate(table_rows):
        for j, val in enumerate(row):
            if j == 0:
                _c(tbl.cell(i+1, j), val, C_NAVY_CELL, C_WHITE, True,
                   PP_ALIGN.CENTER, fs=fs)
            else:
                _c(tbl.cell(i+1, j), val, C_GRAY_CELL, C_BLUE, False,
                   PP_ALIGN.CENTER, fs=fs)


def generate_gdi_core(rows, period_label='', generated_on='', top_n=3,
                      rows_per_slide=12):
    """
    Génère un PPTX reproduisant le rapport « COMITE GESTION DES INCIDENTS –
    Disponibilité et trafic IGW » avec le bloc « Incidents core ».

    Si `top_n` est défini, affiche le TOP `top_n` des incidents les plus
    critiques (durée la plus longue) sur une seule slide. Si `top_n` est None,
    affiche TOUS les incidents, paginés sur plusieurs slides
    (`rows_per_slide` lignes par slide).
    `rows` = liste de dicts : {nature, impact, cause, escalade, duration, duration_sec}
    Retourne un BytesIO.
    """
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    all_rows = rows or []
    if top_n is not None:
        all_rows = all_rows[:top_n]

    def _row(r):
        return [r.get('nature', ''), r.get('impact', ''), r.get('cause', ''),
                r.get('escalade', ''), r.get('duration', '')]

    table_rows = [_row(r) for r in all_rows]

    if not table_rows:
        _gdi_core_slide(prs, [], period_label)
    else:
        chunks = [table_rows[i:i + rows_per_slide]
                  for i in range(0, len(table_rows), rows_per_slide)]
        n_pages = len(chunks)
        for idx, chunk in enumerate(chunks, start=1):
            page_label = f'Page {idx}/{n_pages}' if n_pages > 1 else ''
            _gdi_core_slide(prs, chunk, period_label, page_label)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_igw_dispo(report, top_incidents=None, generated_on=''):
    """
    Génère un PPTX 1 slide « COMITE GESTION DES INCIDENTS – Disponibilité et
    trafic IGW » combinant :
      • à gauche, le tableau « Disponibilité Lien IGW » (Lien / Nombre Inc / Taux),
      • à droite, le « TOP 3 Incidents critiques du mois » (issu de la page Core).
    Retourne un BytesIO.
    """
    from pptx.enum.text import MSO_ANCHOR
    from .igw_dispo import fmt_pct

    links = report.get('links', [])
    month_label = report.get('month_label', '')
    period_label = report.get('period_label', '')
    top_incidents = (top_incidents or [])[:3]

    C_NAVY_CELL = RGBColor(0x0D, 0x24, 0x61)
    C_GRAY_CELL = RGBColor(0xE8, 0xEC, 0xF4)
    C_GREEN = RGBColor(0x22, 0xC5, 0x5E)
    C_ORANGE = RGBColor(0xF5, 0xC0, 0x00)
    C_RED = RGBColor(0xE3, 0x00, 0x13)

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, C_WHITE)

    # Titres
    _txt(sl, 'COMITE GESTION DES INCIDENTS', MARGIN, Inches(0.16),
         Inches(9), Inches(0.45), size=22, bold=True, color=C_BLUE)
    _txt(sl, 'Disponibilité et trafic IGW', MARGIN, Inches(0.60),
         Inches(9), Inches(0.38), size=15, bold=True, color=C_RED)

    # Logo Yas
    try:
        from .gdi_core import yas_logo_bytes
        logo = yas_logo_bytes()
        if logo:
            from PIL import Image as _PILImage
            _img = _PILImage.open(BytesIO(logo))
            ratio = _img.height / _img.width
            lw = Inches(1.2)
            lh = Inches(1.2 * ratio)
            sl.shapes.add_picture(BytesIO(logo), SW - MARGIN - lw, Inches(0.16), lw, lh)
    except Exception:
        pass

    def _cell(cell, val, bg, fg, bold, align, fs=10):
        cell.text = ''
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.06)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = str(val) if val is not None else ''
        r.font.size = Pt(fs)
        r.font.bold = bold
        r.font.color.rgb = fg

    top = Inches(1.55)
    bottom = SH - Inches(0.35)
    gap = Inches(0.30)
    left_w = Inches(6.2)
    right_w = SW - 2 * MARGIN - left_w - gap
    right_x = MARGIN + left_w + gap

    # ── Étiquette gauche ──
    _rect(sl, MARGIN, Inches(1.08), left_w, Inches(0.42), C_YELL)
    _txt(sl, f'Disponibilité Lien IGW {month_label}'.strip(),
         MARGIN, Inches(1.11), left_w, Inches(0.36),
         size=13, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # ── Tableau gauche : Lien / Nombre Inc / Taux ──
    nrows = len(links) + 1
    if nrows < 2:
        nrows = 2
    tblL = sl.shapes.add_table(nrows, 3, MARGIN, top, left_w, bottom - top).table
    tblL.first_row = False
    tblL.horz_banding = False
    lw3 = [3.3, 1.3, 1.6]
    tot = sum(lw3)
    for i, cw in enumerate(lw3):
        tblL.columns[i].width = int(left_w * cw / tot)
    hL = Inches(0.34)
    tblL.rows[0].height = hL
    body_h = int((bottom - top - hL) / max(len(links), 1))
    for i in range(1, nrows):
        tblL.rows[i].height = body_h
    for j, h in enumerate(['Lien', 'Nombre Inc', month_label or 'Taux']):
        _cell(tblL.cell(0, j), h, C_NAVY_CELL, C_WHITE, True, PP_ALIGN.CENTER, fs=11)
    if not links:
        _cell(tblL.cell(1, 0), 'Aucun lien dans le fichier.', C_GRAY_CELL,
              C_BLUE, False, PP_ALIGN.LEFT, fs=10)
        _cell(tblL.cell(1, 1), '', C_GRAY_CELL, C_BLUE, False, PP_ALIGN.CENTER)
        _cell(tblL.cell(1, 2), '', C_GRAY_CELL, C_BLUE, False, PP_ALIGN.CENTER)
    for i, lk in enumerate(links):
        avail = lk['availability']
        is100 = avail >= 99.999
        _cell(tblL.cell(i + 1, 0), lk['short'], C_GRAY_CELL, C_BLUE, True,
              PP_ALIGN.LEFT, fs=10)
        _cell(tblL.cell(i + 1, 1), str(lk['n_inc']), C_GRAY_CELL, C_BLUE, True,
              PP_ALIGN.CENTER, fs=10)
        _cell(tblL.cell(i + 1, 2), fmt_pct(avail),
              C_GREEN if is100 else C_ORANGE,
              C_WHITE if is100 else C_NAVY_CELL, True, PP_ALIGN.CENTER, fs=10)

    # ── Étiquette droite ──
    _rect(sl, right_x, Inches(1.08), right_w, Inches(0.42), C_NAVY_CELL)
    _txt(sl, 'TOP 3 Incidents critiques du mois',
         right_x, Inches(1.11), right_w, Inches(0.36),
         size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ── Tableau droit : TOP 3 ──
    rrows = len(top_incidents) + 1
    if rrows < 2:
        rrows = 2
    rh_total = min(bottom - top, Inches(0.5) + Inches(1.2) * len(top_incidents)) \
        if top_incidents else Inches(2.0)
    tblR = sl.shapes.add_table(rrows, 5, right_x, top, right_w, rh_total).table
    tblR.first_row = False
    tblR.horz_banding = False
    rw3 = [2.6, 1.9, 2.0, 1.4, 1.1]
    totr = sum(rw3)
    for i, cw in enumerate(rw3):
        tblR.columns[i].width = int(right_w * cw / totr)
    hR = Inches(0.40)
    tblR.rows[0].height = hR
    if top_incidents:
        rbody = int((rh_total - hR) / len(top_incidents))
        for i in range(1, rrows):
            tblR.rows[i].height = rbody
    for j, h in enumerate(["Nature de l'incident", 'Impact - Service', 'Cause', 'Escalade', 'Duration']):
        _cell(tblR.cell(0, j), h, C_YELL, C_BLUE, True, PP_ALIGN.LEFT, fs=11)
    if not top_incidents:
        _cell(tblR.cell(1, 0),
              'Importez le fichier de tickets sur la page Core (TOP 3).',
              C_GRAY_CELL, C_BLUE, False, PP_ALIGN.LEFT, fs=10)
        for jj in range(1, 5):
            _cell(tblR.cell(1, jj), '', C_GRAY_CELL, C_BLUE, False, PP_ALIGN.CENTER)
    for i, inc in enumerate(top_incidents):
        _cell(tblR.cell(i + 1, 0), inc.get('nature', ''), C_NAVY_CELL, C_WHITE,
              True, PP_ALIGN.CENTER, fs=10)
        _cell(tblR.cell(i + 1, 1), inc.get('impact', ''), C_GRAY_CELL, C_BLUE,
              False, PP_ALIGN.CENTER, fs=10)
        _cell(tblR.cell(i + 1, 2), inc.get('cause', ''), C_GRAY_CELL, C_BLUE,
              False, PP_ALIGN.CENTER, fs=10)
        _cell(tblR.cell(i + 1, 3), inc.get('escalade', ''), C_GRAY_CELL, C_BLUE,
              False, PP_ALIGN.CENTER, fs=10)
        _cell(tblR.cell(i + 1, 4), inc.get('duration', ''), C_GRAY_CELL, C_BLUE,
              False, PP_ALIGN.CENTER, fs=10)

    if generated_on:
        _txt(sl, f'Généré le {generated_on} — Yas Togo / DT / DOC / iSOC — {period_label}',
             MARGIN, SH - Inches(0.30), SW - 2 * MARGIN, Inches(0.25),
             size=8, color=C_BLUE)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def generate_image_slide(png_buf, generated_on='', footer=''):
    """Place une image PNG (BytesIO) en plein écran sur une diapo unique."""
    from PIL import Image as _PILImage
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    _add_image_slide(prs, png_buf, generated_on, footer)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _add_image_slide(prs, png_buf, generated_on='', footer=''):
    """Ajoute une diapo plein écran contenant l'image PNG à `prs`."""
    from PIL import Image as _PILImage
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, C_WHITE)

    png_buf.seek(0)
    im = _PILImage.open(png_buf)
    w, h = im.size
    png_buf.seek(0)
    ratio = h / w
    avail_w = SW - 2 * MARGIN
    avail_h = SH - 2 * MARGIN - Inches(0.3)
    iw = avail_w
    ih = int(iw * ratio)
    if ih > avail_h:
        ih = int(avail_h)
        iw = int(ih / ratio)
    left = int((SW - iw) / 2)
    top = int((SH - ih - Inches(0.3)) / 2)
    sl.shapes.add_picture(png_buf, left, top, int(iw), int(ih))

    if generated_on or footer:
        _txt(sl, f'Généré le {generated_on} — Yas Togo / DT / DOC / iSOC — {footer}',
             MARGIN, SH - Inches(0.30), SW - 2 * MARGIN, Inches(0.25),
             size=8, color=C_BLUE)
    return sl


def generate_image_deck(png_bufs, generated_on='', footer=''):
    """Regroupe plusieurs images PNG (BytesIO) en un seul PPTX, une par diapo."""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    for png_buf in png_bufs:
        _add_image_slide(prs, png_buf, generated_on, footer)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# RAPPORT NOC TRANSMISSION — DIAPOS NATIVES (MODIFIABLES)
# ═══════════════════════════════════════════════════════════════════════════════

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

C_NAVY_T = RGBColor(0x0D, 0x24, 0x61)
C_RED_T  = RGBColor(0xE3, 0x00, 0x13)
C_GCELL  = RGBColor(0xE8, 0xEC, 0xF4)


def _t_logo(slide):
    from .gdi_core import yas_logo_bytes
    data = yas_logo_bytes()
    if not data:
        return
    try:
        slide.shapes.add_picture(BytesIO(data), SW - Inches(1.55), Inches(0.16),
                                 height=Inches(0.78))
    except Exception:
        pass


def _t_header(slide, subtitle):
    _rect(slide, 0, 0, SW, SH, C_WHITE)
    _txt(slide, 'COMITE GESTION DES INCIDENTS', MARGIN, Inches(0.16),
         Inches(9), Inches(0.45), size=21, bold=True, color=C_BLUE)
    _txt(slide, subtitle, MARGIN, Inches(0.64), Inches(9), Inches(0.35),
         size=14, bold=True, color=C_RED_T)
    _t_logo(slide)


def _t_oval(slide, l, t, w, h, fill, text='', size=13, color=C_WHITE, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    if text:
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return shp


def _t_arrow(slide, l, t, w, h, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _t_box(slide, l, t, w, h, fill, lines, size=9):
    """lines = liste de (texte, couleur, gras)."""
    shp = _rect(slide, l, t, w, h, fill)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, (txt, col, bd) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bd
        r.font.color.rgb = col
    return shp


def _t_cell(cell, text, bg, fg, size=9, bold=False, align=PP_ALIGN.CENTER):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(str(text).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = fg


def _slide_transport_image1(prs, report, generated_on):
    sl = _blank(prs)
    im1 = report['image1']
    _t_header(sl, 'Détails Incident transport')
    _txt(sl, f"{report['total_inc']} incidents enregistrés", MARGIN, Inches(1.12),
         Inches(4), Inches(0.4), size=14, bold=True, color=C_NAVY_T)

    # Cercle central (image officielle si dispo, sinon ovale jaune + texte)
    import os
    cl, ct, cd = Inches(0.55), Inches(2.25), Inches(2.0)
    bg = os.path.join(os.path.dirname(__file__), 'static', 'reports',
                      'backbone_togo.png')
    placed = False
    if os.path.exists(bg):
        try:
            sl.shapes.add_picture(bg, cl, ct, cd, cd)
            placed = True
        except Exception:
            placed = False
    if not placed:
        _t_oval(sl, cl, ct, cd, cd, C_YELL,
                'Backbone de transmission De Yas', size=11, color=C_NAVY_T)

    def branch(node_top, label, blk):
        nx, r = Inches(3.3), Inches(1.0)
        _t_oval(sl, nx, node_top, r, r, C_NAVY_T, label, size=12)
        tx = nx + r + Inches(0.15)
        _txt(sl, f"Inc : {blk['inc']}", tx, node_top + Inches(0.12),
             Inches(2.0), Inches(0.3), size=14, bold=True, color=C_NAVY_T)
        _txt(sl, f"MTTR : {_fmt(blk['mttr_sec'])}", tx, node_top + Inches(0.5),
             Inches(2.3), Inches(0.3), size=14, bold=True, color=C_RED_T)
        _t_arrow(sl, Inches(6.55), node_top + Inches(0.3),
                 Inches(0.75), Inches(0.4), C_NAVY_T)
        bx, bw, bh = Inches(7.5), Inches(2.95), Inches(0.62)
        s, a = blk['sans_impact'], blk['avec_impact']
        _t_box(sl, bx, node_top - Inches(0.05), bw, bh, C_NAVY_T,
               [('INC SANS IMPACT', C_YELL, True),
                (f"Nbre : {s['inc']}    MTTR : {_fmt(s['mttr_sec'])}", C_WHITE, False)])
        _t_box(sl, bx, node_top + Inches(0.63), bw, bh, C_NAVY_T,
               [('INC AVEC IMPACT', C_YELL, True),
                (f"Nbre : {a['inc']}    MTTR : {_fmt(a['mttr_sec'])}", C_WHITE, False)])

    branch(Inches(1.65), 'Backhaul', im1['backhaul'])
    branch(Inches(3.55), 'BackBone', im1['backbone'])

    # Tableau détails Backbone
    details = im1['backbone_details']
    if details:
        _txt(sl, 'Détails incidents BACKBONE', MARGIN, Inches(5.0),
             Inches(5), Inches(0.3), size=12, bold=True, color=C_NAVY_T)
        rows = details[:5]
        widths = [Inches(3.3), Inches(2.4), Inches(2.3)]
        tbl = sl.shapes.add_table(len(rows) + 1, 3, MARGIN, Inches(5.35),
                                  sum(widths, Inches(0)), Inches(0.4) * (len(rows) + 1)).table
        tbl.first_row = False
        tbl.horz_banding = False
        for i, w in enumerate(widths):
            tbl.columns[i].width = w
        for i, h in enumerate(["Nature de l'incident", 'Impact - Service', 'Cause']):
            _t_cell(tbl.cell(0, i), h, C_BLUE3, C_YELL, size=10, bold=True)
        for ri, row in enumerate(rows, start=1):
            _t_cell(tbl.cell(ri, 0), row['lien'], C_NAVY_T, C_WHITE,
                    size=9, align=PP_ALIGN.LEFT)
            _t_cell(tbl.cell(ri, 1), row['impact'], C_GCELL, C_NAVY_T, size=9)
            _t_cell(tbl.cell(ri, 2), row['cause'], C_GCELL, C_NAVY_T, size=9)

    if generated_on:
        _txt(sl, f'Généré le {generated_on} — Yas Togo / DT / DOC / iSOC — '
             f"{report.get('period_label', '')}", MARGIN, SH - Inches(0.32),
             Inches(11), Inches(0.25), size=8, color=C_BLUE)
    return sl


def _slide_transport_image2(prs, report, generated_on):
    sl = _blank(prs)
    im2 = report['image2']
    regions = [r for r in im2['regions'] if r['canonical'] or r['has_data']]
    _t_header(sl, 'Count Inc & MTTR par Métier et par Régions')

    band, wm, wi, wt = Inches(1.25), Inches(1.65), Inches(0.85), Inches(1.45)
    tw = band + wm + wi + wt
    row_h = Inches(0.26)
    cols_x = [MARGIN, MARGIN + tw + Inches(0.4)]

    def draw_region(x, y, reg):
        nrows = 1 + len(reg['metiers'])
        tbl = sl.shapes.add_table(nrows, 4, x, y, tw, row_h * nrows).table
        tbl.first_row = False
        tbl.horz_banding = False
        for i, w in enumerate([band, wm, wi, wt]):
            tbl.columns[i].width = w
        for i in range(nrows):
            tbl.rows[i].height = row_h
        # bandeau région (col 0 fusionnée)
        tbl.cell(0, 0).merge(tbl.cell(nrows - 1, 0))
        _t_cell(tbl.cell(0, 0), reg['region'], C_YELL, C_NAVY_T, size=11, bold=True)
        for i, h in enumerate(['Métier', 'Inc', 'MTTR'], start=1):
            _t_cell(tbl.cell(0, i), h, C_MGRAY, C_NAVY_T, size=9, bold=True)
        for ri, m in enumerate(reg['metiers'], start=1):
            hot = m['mttr_sec'] >= 5 * 3600
            mttr = _fmt(m['mttr_sec']) if m['inc'] else '0:00:00'
            _t_cell(tbl.cell(ri, 1), m['metier'], C_WHITE, C_NAVY_T,
                    size=9, align=PP_ALIGN.LEFT)
            _t_cell(tbl.cell(ri, 2), str(m['inc']), C_WHITE, C_NAVY_T, size=9, bold=True)
            _t_cell(tbl.cell(ri, 3), mttr, C_RED_T if hot else C_WHITE,
                    C_WHITE if hot else C_NAVY_T, size=9, bold=True)
        return nrows * row_h + Inches(0.18)

    n = len(regions)
    columns = [regions[:(n + 1) // 2], regions[(n + 1) // 2:]]
    ys = [Inches(1.3), Inches(1.3)]
    for ci, lst in enumerate(columns):
        for reg in lst:
            ys[ci] += draw_region(cols_x[ci], ys[ci], reg)

    # Box BACKBONE DWDM (sous la colonne la moins remplie)
    dwdm = im2['backbone_dwdm']
    ci = 0 if ys[0] <= ys[1] else 1
    by = ys[ci]
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, cols_x[ci], by, tw, Inches(1.0))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C_NAVY_T
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate([
            f"{dwdm['count']} Indisponibilité(s) du BACKBONE DWDM",
            f"MTTR : {_fmt(dwdm['mttr_sec'])}", dwdm['services']]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(10)
        r.font.bold = (i == 0)
        r.font.color.rgb = C_WHITE

    if generated_on:
        _txt(sl, f'Généré le {generated_on} — Yas Togo / DT / DOC / iSOC — '
             f"{report.get('period_label', '')}", MARGIN, SH - Inches(0.32),
             Inches(11), Inches(0.25), size=8, color=C_BLUE)
    return sl


def _slide_transport_image3(prs, report, generated_on):
    sl = _blank(prs)
    clients = report['image3']['clients']
    _t_header(sl, 'Disponibilité clients IPT et IPLC')
    band = _rect(sl, MARGIN, Inches(1.08), Inches(2.3), Inches(0.38), C_NAVY_T)
    _txt(sl, 'Clients IPT & IPLC', MARGIN + Inches(0.1), Inches(1.12),
         Inches(2.1), Inches(0.3), size=12, bold=True, color=C_WHITE)

    widths = [Inches(3.2), Inches(1.3), Inches(1.9), Inches(2.4)]
    tw = sum(widths, Inches(0))
    left = (SW - tw) / 2
    nrows = 1 + len(clients)
    top = Inches(1.6)
    avail = SH - top - Inches(0.35)
    tbl = sl.shapes.add_table(nrows, 4, left, top, tw, avail).table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, w in enumerate(widths):
        tbl.columns[i].width = w
    for i, h in enumerate(['LIENS', 'Nbre Inc', 'Durée', 'TAUX DE DISPONIBILITE']):
        _t_cell(tbl.cell(0, i), h, C_MGRAY, C_NAVY_T, size=11, bold=True)
    for ri, c in enumerate(clients, start=1):
        nb_bg = C_YELL if c['inc'] > 0 else C_WHITE
        durs = '\n'.join(_fmt(x) for x in c['durations']) if c['durations'] else '0:00:00'
        taux = f"{c['taux']:.2f}".replace('.', ',') + ' %'
        _t_cell(tbl.cell(ri, 0), c['name'], C_WHITE, C_NAVY_T, size=11, bold=True)
        _t_cell(tbl.cell(ri, 1), str(c['inc']), nb_bg,
                C_RED_T if c['inc'] > 0 else C_NAVY_T, size=11, bold=True)
        _t_cell(tbl.cell(ri, 2), durs, C_WHITE, C_NAVY_T, size=10)
        _t_cell(tbl.cell(ri, 3), taux, C_WHITE, C_NAVY_T, size=11, bold=True)

    if generated_on:
        _txt(sl, f'Généré le {generated_on} — Yas Togo / DT / DOC / iSOC — '
             f"{report.get('period_label', '')}", MARGIN, SH - Inches(0.32),
             Inches(11), Inches(0.25), size=8, color=C_BLUE)
    return sl


_TRANSPORT_SLIDES = {
    'image1': _slide_transport_image1,
    'image2': _slide_transport_image2,
    'image3': _slide_transport_image3,
}


def generate_transport_editable(report, generated_on='',
                                images=('image1', 'image2', 'image3')):
    """Rapport NOC transmission en diapos natives PowerPoint (modifiables)."""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    for im in images:
        fn = _TRANSPORT_SLIDES.get(im)
        if fn:
            fn(prs, report, generated_on)
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# EXPORT STATISTIQUES PPTX
# ═══════════════════════════════════════════════════════════════════════════════

def _cover_stats(prs, period_label, generated_on):
    sl = _blank(prs)
    _rect(sl, 0, 0, SW, SH, C_BLUE)
    _rect(sl, 0, 0, SW, Inches(0.12), C_YELL)
    _rect(sl, 0, SH - Inches(0.12), SW, Inches(0.12), C_YELL)
    _rect(sl, SW - Inches(0.12), 0, Inches(0.12), SH, C_YELL)
    _rect(sl, 0, 0, Inches(0.12), SH, C_YELL)
    _txt(sl, 'Yas Togo / DT / DOC / iSOC', MARGIN, Inches(0.5),
         SW - 2*MARGIN, Inches(0.4), size=11,
         color=RGBColor(0x80, 0xA0, 0xD0), align=PP_ALIGN.CENTER)
    _txt(sl, 'STATISTIQUES', MARGIN, Inches(1.5),
         SW - 2*MARGIN, Inches(1.1), size=60, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, 'RÉSEAU MOBILE', MARGIN, Inches(2.5),
         SW - 2*MARGIN, Inches(0.9), size=40, bold=True, color=C_YELL, align=PP_ALIGN.CENTER)
    _txt(sl, period_label.upper() if period_label else '', MARGIN, Inches(3.7),
         SW - 2*MARGIN, Inches(0.8), size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _txt(sl, f'Généré le {generated_on}', MARGIN, SH - Inches(0.8),
         SW - 2*MARGIN, Inches(0.35), size=10,
         color=RGBColor(0x80, 0xA0, 0xD0), align=PP_ALIGN.CENTER)


def _slide_stats_kpis(prs, total_incidents, total_outage_sec,
                      total_duree_sec, nb_escalades, mois_label):
    sl = _blank(prs)
    _header(sl, "Vue d'ensemble — Réseau Mobile", mois_label=mois_label)

    def _fh(s):
        return f'{s/3600:.1f}h' if s else '0.0h'

    mttr_sec = total_duree_sec / total_incidents if total_incidents else 0
    kpis = [
        ('Total Incidents',   str(total_incidents),  'PÉRIODE',    C_BLUE),
        ('Outage Total',      _fh(total_outage_sec), 'RÉSEAU',     RGBColor(0xD9, 0x50, 0x00)),
        ('MTTR Moyen',        _fh(mttr_sec),         'PAR INCIDENT', RGBColor(0xC9, 0x77, 0x00)),
        ('Escalades Actives', str(nb_escalades),     'MÉTIERS',    RGBColor(0x05, 0x80, 0x50)),
    ]
    _kpi_bar(sl, kpis)


def _slide_stats_escalades(prs, escalades_sorted, total_incidents,
                            total_outage_sec, total_duree_sec, mois_label):
    sl = _blank(prs)
    _header(sl, 'Synthèse par Escalade', mois_label=mois_label)

    rows, fmts = [], {}
    for i, (esc, v) in enumerate(escalades_sorted):
        count = v['count']
        if not count:
            continue
        mttr_sec = v['duree_sec'] / count
        pct = f"{round(v['outage_sec'] / total_outage_sec * 100)}%" if total_outage_sec else '0%'
        rows.append([esc, count, _fmt(v['duree_sec']), _fmt(mttr_sec), _fmt(v['outage_sec']), pct])
        ri = len(rows) - 1
        if count >= 20:
            fmts[(ri, 1)] = (C_RED_BG, C_RED_FG)
        elif count >= 10:
            fmts[(ri, 1)] = (C_YELL_BG, C_YELL_FG)
        else:
            fmts[(ri, 1)] = (C_GREEN_BG, C_GREEN_FG)

    total_mttr = total_duree_sec / total_incidents if total_incidents else 0
    rows.append(['TOTAL', total_incidents, _fmt(total_duree_sec),
                 _fmt(total_mttr), _fmt(total_outage_sec), '100%'])
    last = len(rows) - 1
    for j in range(6):
        fmts[(last, j)] = (C_BLUE3, C_WHITE)

    _table(sl, ['ESCALADE', 'INC', 'DURÉE', 'MTTR', 'OUTAGE', '% OUTAGE'],
           rows, col_widths=[4.5, 1, 2, 2, 2, 1.5],
           cell_fmts=fmts, font_size=9)


def _slide_stats_degraded(prs, degraded_top10, site_top_cause, mois_label):
    sl = _blank(prs)
    _header(sl, 'Sites les Plus Dégradés — Durée Outage', mois_label=mois_label)

    max_dur = degraded_top10[0][1] if degraded_top10 else 1
    rows, fmts = [], {}
    for i, (site, dur_sec) in enumerate(degraded_top10):
        cause = (site_top_cause or {}).get(site) or 'N/A'
        rows.append([str(i + 1), site[:35], f'{dur_sec/3600:.1f}h', cause[:45]])
        ri = len(rows) - 1
        if dur_sec >= max_dur * 0.7:
            fmts[(ri, 2)] = (C_RED_BG, C_RED_FG)
        elif dur_sec >= max_dur * 0.4:
            fmts[(ri, 2)] = (C_YELL_BG, C_YELL_FG)
        else:
            fmts[(ri, 2)] = (C_GREEN_BG, C_GREEN_FG)

    _table(sl, ['#', 'SITE', 'DURÉE OUTAGE', 'CAUSE PRINCIPALE'],
           rows, col_widths=[0.5, 4, 2, 6.5], cell_fmts=fmts, font_size=10)


def _slide_stats_recurrence(prs, sites_top10, mois_label):
    sl = _blank(prs)
    _header(sl, 'Récurrence des Sites — Top Incidents', mois_label=mois_label)

    max_cnt = sites_top10[0][1] if sites_top10 else 1
    rows, fmts = [], {}
    for i, (site, count) in enumerate(sites_top10):
        rows.append([str(i + 1), site[:55], count])
        ri = len(rows) - 1
        if count >= max_cnt * 0.7:
            fmts[(ri, 2)] = (C_RED_BG, C_RED_FG)
        elif count >= max_cnt * 0.4:
            fmts[(ri, 2)] = (C_YELL_BG, C_YELL_FG)
        else:
            fmts[(ri, 2)] = (C_GREEN_BG, C_GREEN_FG)

    _table(sl, ['#', 'SITE', 'NB INCIDENTS'],
           rows, col_widths=[0.5, 9, 2], cell_fmts=fmts, font_size=10)


def _slide_stats_causes(prs, causes_dur_top10, mois_label):
    sl = _blank(prs)
    _header(sl, "Top Causes — Par Durée d'Outage", mois_label=mois_label)

    total_dur = sum(d for _, d in causes_dur_top10) or 1
    rows = []
    for i, (cause, dur_sec) in enumerate(causes_dur_top10):
        pct = round(dur_sec / total_dur * 100)
        rows.append([str(i + 1), cause[:60], f'{dur_sec/3600:.1f}h', f'{pct}%'])

    _table(sl, ['#', 'CAUSE', 'DURÉE OUTAGE', '% DU TOTAL'],
           rows, col_widths=[0.5, 8.5, 2, 1.5], font_size=10)


def _slide_stats_dispo(prs, dispo_table, semaine_labels, mois_label):
    """Slide disponibilité : résumé min/moy par escalade."""
    sl = _blank(prs)
    _header(sl, 'Disponibilité Réseau — Taux par Équipement', mois_label=mois_label)

    rows, fmts = [], {}
    for esc, periods in dispo_table.items():
        vals = [v for v in periods.values() if v is not None]
        if not vals:
            continue
        mn  = min(vals)
        moy = sum(vals) / len(vals)
        derniere_lbl = semaine_labels[-1] if semaine_labels else ''
        derniere = periods.get(derniere_lbl)
        rows.append([esc,
                     f'{mn:.2f}%',
                     f'{moy:.2f}%',
                     f'{derniere:.2f}%' if derniere is not None else '—'])
        ri = len(rows) - 1
        val_ref = mn
        if val_ref < 99.0:
            fmts[(ri, 1)] = (C_RED_BG, C_RED_FG)
        elif val_ref < 99.5:
            fmts[(ri, 1)] = (C_YELL_BG, C_YELL_FG)
        else:
            fmts[(ri, 1)] = (C_GREEN_BG, C_GREEN_FG)

    if rows:
        _table(sl, ['ÉQUIPEMENT', 'MIN (%)', 'MOY (%)', 'DERNIÈRE PÉRIODE'],
               rows, col_widths=[4, 2.5, 2.5, 3.5],
               cell_fmts=fmts, font_size=10)
    else:
        _txt(sl, 'Aucune donnée de disponibilité pour cette période.',
             MARGIN, CONTENT_TOP + Inches(1), SW - 2*MARGIN, Inches(0.5),
             size=14, color=C_BLUE3)


def generate_statistiques_pptx(
    escalades_sorted,
    total_incidents,
    total_outage_sec,
    total_duree_sec,
    degraded_top10,
    site_top_cause,
    sites_top10,
    causes_dur_top10,
    dispo_table,
    semaine_labels,
    period_label='',
    generated_on='',
):
    """
    Génère le PPTX export des statistiques réseau mobile.
    Retourne un BytesIO prêt à être envoyé en réponse HTTP.
    """
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    mois_label = period_label or ''

    _cover_stats(prs, period_label, generated_on)

    nb_actives = sum(1 for _, v in escalades_sorted if v['count'] > 0)
    _slide_stats_kpis(prs, total_incidents, total_outage_sec,
                      total_duree_sec, nb_actives, mois_label)

    if escalades_sorted:
        _slide_stats_escalades(prs, escalades_sorted, total_incidents,
                                total_outage_sec, total_duree_sec, mois_label)

    if degraded_top10:
        _slide_stats_degraded(prs, degraded_top10, site_top_cause, mois_label)

    if sites_top10:
        _slide_stats_recurrence(prs, sites_top10, mois_label)

    if causes_dur_top10:
        _slide_stats_causes(prs, causes_dur_top10, mois_label)

    if dispo_table:
        _slide_stats_dispo(prs, dispo_table, semaine_labels, mois_label)

    _closing(prs)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

